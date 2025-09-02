from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw

# Paths
img_path = r"D:\TLDP_Project\image\images (3).jpg"
cropped_path = "circle_image.png"

# Make circular crop with Pillow
im = Image.open(img_path).convert("RGBA")
bigsize = (im.size[0] * 3, im.size[1] * 3)
mask = Image.new("L", bigsize, 0)
draw = ImageDraw.Draw(mask)
draw.ellipse((0, 0) + bigsize, fill=255)
mask = mask.resize(im.size, Image.LANCZOS)
im.putalpha(mask)
im.save(cropped_path, format="PNG")

# Insert into PPT
prs = Presentation("test.pptx")
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text") and "{{image1}}" in shape.text:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(cropped_path, left, top, width, height)

prs.save("output.pptx")
