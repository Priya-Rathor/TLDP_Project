from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta
from PIL import Image, ImageDraw
import platform
import os
import requests
import re
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os, re
# -------------------------
# Layout processing for slide 3
# -------------------------
def process_layout_slide(prs, slide, layout_img_path, bg_img_path):
    """Process slide 3 with layout image and background."""
    slide_width, slide_height = prs.slide_width, prs.slide_height

    # Replace {{Layout1}} placeholder
    layout_replaced = False
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and (("{{Layout1}}" in shape.text) or ("{{layout1}}" in shape.text)):
            print("🔄 Replacing {{Layout1}} placeholder")
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            if layout_img_path and os.path.exists(layout_img_path):
                slide.shapes.add_picture(layout_img_path, left, top, width, height)
                print("✅ Layout image added")
                layout_replaced = True
            break

    # Add background image
    if bg_img_path and os.path.exists(bg_img_path):
        pic = slide.shapes.add_picture(bg_img_path, 0, 0, slide_width, slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        print("✅ Layout background image added")

    return slide


# -------------------------
# Constants
# -------------------------
if platform.system() == "Windows":
    DAY_FORMAT = "%#d-%b"
else:
    DAY_FORMAT = "%-d-%b"


# -------------------------
# Helper: Load image from URL or path
# -------------------------
def get_local_image(path_or_url, tmp_name="temp_img.png"):
    """Download image from URL or validate local path."""
    if not path_or_url:
        return None
    
    if isinstance(path_or_url, str) and path_or_url.startswith("http"):
        try:
            resp = requests.get(path_or_url, timeout=10)
            resp.raise_for_status()
            with open(tmp_name, "wb") as f:
                f.write(resp.content)
            print(f"✅ Downloaded: {tmp_name}")
            return tmp_name
        except Exception as e:
            print(f"⚠️ Failed to download image {path_or_url}: {e}")
            return None
    
    return path_or_url if os.path.exists(path_or_url) else None


# -------------------------
# Helper: Make circular crop
# -------------------------
def make_circle_image(img_path, output_path="circle_image.png"):
    """Create a circular cropped version of the image."""
    im = Image.open(img_path).convert("RGBA")
    bigsize = (im.size[0] * 3, im.size[1] * 3)
    mask = Image.new("L", bigsize, 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0) + bigsize, fill=255)
    mask = mask.resize(im.size, Image.LANCZOS)
    im.putalpha(mask)
    im.save(output_path, format="PNG")
    return output_path


# -------------------------
# SIMPLE DIAGNOSTIC: Show all text in slides
# -------------------------
def show_slide_text(prs):
    """Show all text content in slides to find placeholders."""
    print("\n🔍 CHECKING ALL SLIDES FOR PLACEHOLDERS:")
    print("=" * 50)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n📄 SLIDE {slide_num}:")
        
        found_placeholders = []
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        print(f"  Text: '{text}'")
                        # Look for any placeholder patterns
                        if "{{" in text and "}}" in text:
                            placeholders = re.findall(r"\{\{[^}]+\}\}", text)
                            if placeholders:
                                found_placeholders.extend(placeholders)
                                print(f"    🎯 PLACEHOLDERS: {placeholders}")
            except:
                pass
        
        if not found_placeholders:
            print("  (No placeholders found)")
    
    print("=" * 50)


# -------------------------
# FIXED: Replace ALL image placeholders in ALL slides
# -------------------------


def replace_all_image_placeholders(prs, extra_images):
    """
    Replace ALL {{Image1}}, {{Image2}}, etc. across ALL slides,
    inserting images at their ORIGINAL size (scaled only if too big).
    """
    print(f"\n🖼️ REPLACING IMAGE PLACEHOLDERS IN ALL SLIDES")
    print(f"📋 Available images: {len([img for img in extra_images if img])}")
    
    total_replaced = 0
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n🔍 Checking Slide {slide_num}...")

        shapes_to_check = list(slide.shapes)
        
        for shape in shapes_to_check:
            try:
                if not hasattr(shape, "text") or not shape.text:
                    continue

                # Look for {{ImageX}} placeholders
                matches = re.findall(r"\{\{[Ii]mage(\d+)\}\}", shape.text)
                
                if matches:
                    for match in matches:
                        placeholder_num = int(match)
                        img_index = placeholder_num - 1  # zero-based index
                        
                        print(f"  🎯 Found {{Image{placeholder_num}}} placeholder")

                        # Record placeholder position
                        left, top = shape.left, shape.top
                        
                        # Remove placeholder
                        slide.shapes._spTree.remove(shape._element)
                        print("  🗑️ Removed placeholder shape")
                        
                        # Insert image if available
                        if 0 <= img_index < len(extra_images) and extra_images[img_index]:
                            img_path = extra_images[img_index]
                            if os.path.exists(img_path):
                                # Get original image size
                                with Image.open(img_path) as img:
                                    px_width, px_height = img.size
                                
                                # Convert px → EMU (1 inch = 96 px = 914400 EMU)
                                emu_width = int(px_width * 914400 / 96)
                                emu_height = int(px_height * 914400 / 96)

                                # Scale down if too big for slide
                                if emu_width > slide_width or emu_height > slide_height:
                                    scale = min(slide_width / emu_width, slide_height / emu_height)
                                    emu_width = int(emu_width * scale)
                                    emu_height = int(emu_height * scale)

                                # Add picture at original/native size
                                slide.shapes.add_picture(img_path, left, top, width=emu_width, height=emu_height)
                                print(f"  ✅ Added original-size image: {os.path.basename(img_path)}")
                                total_replaced += 1
                            else:
                                print(f"  ❌ Image not found: {img_path}")
                        else:
                            print(f"  ⚠️ No image available for placeholder {placeholder_num}")
                        
                        break  # handle one match per shape

            except Exception as e:
                print(f"  ❌ Error processing shape: {e}")
    
    print(f"\n📊 TOTAL IMAGES REPLACED: {total_replaced}")
    return prs

# -------------------------
# Calendar helpers
# -------------------------
def build_mapping(start_date=None):
    """Build date mapping for calendar placeholders."""
    if start_date is None:
        start_date = datetime.today()
    
    mapping = {}
    for i in range(1, 8):
        date = start_date + timedelta(days=i - 1)
        mapping[f"{{{{day{i}}}}}"] = date.strftime("%a")
    
    for i in range(1, 50):
        date = start_date + timedelta(days=i - 1)
        mapping[f"{{{{d{i}}}}}"] = date.strftime(DAY_FORMAT)
    
    return mapping


def replace_text_in_frame(text_frame, mapping):
    """Replace text placeholders in a text frame."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            for ph, val in mapping.items():
                if ph in run.text:
                    run.text = run.text.replace(ph, val)


def iter_shapes(shapes):
    """Recursively iterate through all shapes including groups."""
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # group shape
            yield from iter_shapes(shp.shapes)


def update_calendar_with_bg(prs, slide, image_path, start_date=None):
    """Update calendar with background image and date mappings."""
    mapping = build_mapping(start_date)
    slide_width, slide_height = prs.slide_width, prs.slide_height

    if image_path and os.path.exists(image_path):
        pic = slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        print("✅ Calendar background added")

    for shp in iter_shapes(slide.shapes):
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    replace_text_in_frame(cell.text_frame, mapping)
        elif getattr(shp, "has_text_frame", False):
            replace_text_in_frame(shp.text_frame, mapping)
    
    return slide


def replace_text_in_ppt(slide, text_map):
    """Replace text placeholders in slide."""
    for shp in iter_shapes(slide.shapes):
        if getattr(shp, "has_text_frame", False):
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    for ph, val in text_map.items():
                        if ph in run.text:
                            run.text = run.text.replace(ph, val)
    return slide


def replace_with_circle_image(slide, img_path):
    cropped_path = None
    if not img_path or not os.path.exists(img_path):
        print(f"⚠️ Circle image not found: {img_path}")
        return slide, None
    
    cropped_path = make_circle_image(img_path)
    
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and (("{{Image1}}" in shape.text) or ("{{image1}}" in shape.text)):
            print(f"🔄 Replacing {{Image1}} with circular image")
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(cropped_path, left, top, width, height)
            print("✅ Circle image added")
            break
    
    return slide, cropped_path



def cleanup_temp_files(files):
    """Clean up temporary files."""
    cleaned = 0
    for f in files:
        try:
            if f and os.path.exists(f):
                os.remove(f)
                cleaned += 1
        except Exception as e:
            print(f"⚠️ Could not delete {f}: {e}")
    print(f"🧹 Cleaned up {cleaned} temporary files")


# -------------------------
# MAIN: Create Brochure (SIMPLIFIED)
# -------------------------
def create_brochure_ppt(template_path, output_path, form_data,
                        circle_img, calendar_bg, layout_img, layout_bg, extra_images):
    """
    Main function to create brochure PPT.
    Handles Slide 1 (text + circle), Slide 2 (calendar), Slide 3 (layout + background),
    and replaces all {{ImageX}} placeholders across all slides.
    """
    print(f"🔄 Creating brochure PPT from {template_path}")
    
    # Load presentation
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    prs = Presentation(template_path)
    print(f"📊 Loaded presentation with {len(prs.slides)} slides")

    # Show what's in the slides (debug helper)
    show_slide_text(prs)

    # -------------------------
    # Process images
    # -------------------------
    print("\n🔄 Processing input images...")
    temp_files = []
    
    # Extra images
    processed_extra_images = []
    for i, img in enumerate(extra_images):
        if img:
            local_img = get_local_image(img, f"extra_{i}.png")
            if local_img:
                processed_extra_images.append(local_img)
                temp_files.append(local_img)
                print(f"✅ Processed image {i+1}: {os.path.basename(local_img)}")
            else:
                processed_extra_images.append(None)
                print(f"❌ Failed to process image {i+1}")
        else:
            processed_extra_images.append(None)
    
    print(f"📋 Total processed images: {len([img for img in processed_extra_images if img])}")

    # -------------------------
    # Slide 1 (Text + Circle Image)
    # -------------------------
    circle_img = get_local_image(circle_img, "circle.png")
    if circle_img: 
        temp_files.append(circle_img)
        if len(prs.slides) > 0:
            print("\n🔄 Processing Slide 1...")
            slide = prs.slides[0]
            text_map = {
                "Project Name": form_data.get("Project Name", ""),
                "Project Type": form_data.get("What is the nature of your project?", ""),
                "space To be Designed": form_data.get("Space(S) to be designed", ""),
                "Room size": form_data.get("What is the area size?", ""),
                "style(s) selected": form_data.get("Which style(s) do you like?", ""),
                "Location": f"{form_data.get('City', '')}, {form_data.get('Country', '')}",
            }
            replace_text_in_ppt(slide, {k: v for k, v in text_map.items() if v})
            slide, cropped_circle = replace_with_circle_image(slide, circle_img)
            if cropped_circle:
                temp_files.append(cropped_circle)  # cleanup later

    # -------------------------
    # Slide 2 (Calendar)
    # -------------------------
    calendar_bg = get_local_image(calendar_bg, "calendar_bg.png")
    if calendar_bg: 
        temp_files.append(calendar_bg)
        if len(prs.slides) > 1:
            print("\n🔄 Processing Slide 2 (Calendar)...")
            slide = prs.slides[1]
            update_calendar_with_bg(prs, slide, calendar_bg)
     
    # -------------------------
    # Slide 3 (Layout + Background)
    # -------------------------
    layout_img = get_local_image(layout_img, "layout_img.png")
    layout_bg = get_local_image(layout_bg, "layout_bg.png")
    if layout_img: temp_files.append(layout_img)
    if layout_bg: temp_files.append(layout_bg)

    if len(prs.slides) > 2:
        print("\n🔄 Processing Slide 3 (Layout + Background)...")
        slide = prs.slides[2]
        process_layout_slide(prs, slide, layout_img, layout_bg)
       
    # -------------------------
    # Replace ALL {{ImageX}} placeholders in ALL slides
    # -------------------------
    replace_all_image_placeholders(prs, processed_extra_images)

    # -------------------------
    # Save final PPT
    # -------------------------
    print(f"\n💾 Saving presentation to: {output_path}")
    prs.save(output_path)
    print(f"✅ Brochure PPT created: {output_path}")

    # -------------------------
    # Cleanup temp files
    # -------------------------
    cleanup_temp_files(temp_files)

    return output_path
