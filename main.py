from fastapi import FastAPI, Request
from email_utils import send_email_with_ppt
from fastapi.responses import JSONResponse
from pptx import Presentation
import os, re, json, requests
from io import BytesIO
from docx import Document
import fitz  # PyMuPDF
from zipfile import ZipFile
from pptx.util import Inches
from brochure import create_brochure_ppt
from PIL import Image
from city import insert_city_image_in_ppt
from style import filter_ppt

app = FastAPI()

TEMPLATE_PATH = "template.pptx"
OUTPUT_PATH = "output.pptx"
MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "eyJhbGciOiJIUzI1NiJ9.eyJ0aWQiOjU0NjI5MjM1NywiYWFpIjoxMSwidWlkIjo3NDc3Njk5NywiaWFkIjoiMjAyNS0wOC0wNFQwOTo0MzowNS4wMDBaIiwicGVyIjoibWU6d3JpdGUiLCJhY3RpZCI6MTIxNDMyMDQsInJnbiI6InVzZTEifQ.yYeelRXHOZlaxwYHBAvi6eXRzD2fNn1H-jX-Pd8Ukcw")
MONDAY_API_URL = "https://api.monday.com/v2"

# Track processed items to prevent duplicate emails
PROCESSED_ITEMS = set()
EMAIL_SENT_LOG = "email_sent_log.txt"

def load_processed_items():
    """Load previously processed item IDs from file"""
    global PROCESSED_ITEMS
    if os.path.exists(EMAIL_SENT_LOG):
        try:
            with open(EMAIL_SENT_LOG, 'r') as f:
                PROCESSED_ITEMS = set(line.strip() for line in f if line.strip())
            print(f"📋 Loaded {len(PROCESSED_ITEMS)} previously processed items")
        except Exception as e:
            print(f"⚠️ Error loading processed items: {e}")
            PROCESSED_ITEMS = set()
    else:
        PROCESSED_ITEMS = set()

def mark_item_as_processed(item_id):
    """Mark an item as processed and save to file"""
    global PROCESSED_ITEMS
    PROCESSED_ITEMS.add(str(item_id))
    try:
        with open(EMAIL_SENT_LOG, 'a') as f:
            f.write(f"{item_id}\n")
        print(f"✅ Marked item {item_id} as processed")
    except Exception as e:
        print(f"⚠️ Error saving processed item: {e}")

def is_item_processed(item_id):
    """Check if an item has already been processed"""
    return str(item_id) in PROCESSED_ITEMS

# Load processed items on startup
load_processed_items()

def get_image_dimensions(img_path_or_bytes):
    """Get original dimensions of an image. Returns (width, height) in pixels or None if failed."""
    try:
        if isinstance(img_path_or_bytes, str) and img_path_or_bytes.startswith("http"):
            response = requests.get(img_path_or_bytes, timeout=10)
            response.raise_for_status()
            img_bytes = BytesIO(response.content)
            img = Image.open(img_bytes)
        elif isinstance(img_path_or_bytes, str):
            img = Image.open(img_path_or_bytes)
        else:
            img = Image.open(img_path_or_bytes)
        
        return img.size  # Returns (width, height)
    except Exception as e:
        print(f"⚠️ Failed to get image dimensions: {e}")
        return None

def calculate_image_size_for_slide_fixed_height(img_width, img_height, placeholder_width, placeholder_height, fixed_height_px=8000):
    """
    Calculate image size with FIXED height of 540px and proportional width.
    
    Args:
        img_width, img_height: Original image dimensions in pixels
        placeholder_width, placeholder_height: Placeholder dimensions in PowerPoint EMUs
        fixed_height_px: Fixed height in pixels (default 540)
    
    Returns:
        (new_width, new_height) in PowerPoint units (EMUs)
    """
    print(f"🔍 Input: img={img_width}x{img_height}px, placeholder={placeholder_width}x{placeholder_height}EMUs, fixed_height={fixed_height_px}px")
    
    # Convert pixels to PowerPoint EMUs (English Metric Units)
    # 1 inch = 914400 EMUs, assuming 72 DPI (standard for most images)
    PIXELS_PER_INCH = 72
    EMUS_PER_INCH = 914400
    
    # Calculate aspect ratio (width/height)
    aspect_ratio = img_width / img_height if img_height > 0 else 1.0
    print(f"📐 Original aspect ratio (w/h): {aspect_ratio:.3f}")
    
    # Set fixed height and calculate proportional width
    new_height_px = fixed_height_px
    new_width_px = int(fixed_height_px * aspect_ratio)
    
    print(f"📏 Fixed dimensions: {new_width_px}x{new_height_px}px")
    
    # Convert to EMUs
    new_width_emu = int((new_width_px / PIXELS_PER_INCH) * EMUS_PER_INCH)
    new_height_emu = int((new_height_px / PIXELS_PER_INCH) * EMUS_PER_INCH)
    
    # Check if image fits within placeholder bounds
    if new_width_emu > placeholder_width:
        scale_factor = placeholder_width / new_width_emu
        new_width_emu = placeholder_width
        new_height_emu = int(new_height_emu * scale_factor)
        final_height_px = int(new_height_emu / EMUS_PER_INCH * PIXELS_PER_INCH)
        print(f"⚠️ Width exceeded placeholder, scaled down to: {new_width_px}x{final_height_px}px")
    
    print(f"📏 Final EMU dimensions: {new_width_emu}x{new_height_emu}")
    print(f"📊 Final pixel equivalent: {new_width_emu/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}x{new_height_emu/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}px")
    
    return new_width_emu, new_height_emu

def replace_placeholders_with_images(pptx_path, output_path, categorized_images):
    """
    Replace placeholders like {{Image1}}, {{Layout2}}, {{Elevation1}}, {{Inspiration1}} 
    with categorized images positioned at the placeholder's exact location.
    
    Features:
    - Fixed 540px HEIGHT for all images
    - Universal image looping: ALL categories use ALL available images when needed
    - Keep placeholder only if NO images exist across ALL categories
    """
    prs = Presentation(pptx_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"📏 Slide dimensions: {slide_width} x {slide_height}")

    # Build category mapping
    category_map = {
        "Layout": categorized_images.get("floor_plans", []),
        "Elevation": categorized_images.get("elevation_drawings", []),
        "Inspiration": categorized_images.get("inspiration_images", []),
        "Image": categorized_images.get("existing_pictures", []),
    }

    # Create a universal image pool from ALL categories
    all_images = []
    for category, images in category_map.items():
        all_images.extend(images)
    
    print("📁 Available images per category:")
    for category, images in category_map.items():
        print(f"   {category}: {len(images)} images")
    print(f"🌐 UNIVERSAL POOL: {len(all_images)} total images available for looping")

    slides_to_delete = []

    for slide_idx, slide in enumerate(prs.slides):
        replaced_any = False
        found_placeholder = False

        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            if shape.has_text_frame:
                text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs).strip()
            else:
                text = shape.text.strip()

            # Match placeholders like {{CategoryN}} but EXCLUDE style placeholders
            match = re.match(r"\{\{(\w+)(\d+)\}\}", text)
            if not match:
                continue

            category, num = match.group(1), int(match.group(2)) - 1
            
            # Skip style placeholders (they're handled separately)
            if category.lower() == "style":
                continue

            found_placeholder = True
            
            # Universal image selection strategy
            category_images = category_map.get(category, [])
            selected_image = None
            
            # Strategy 1: Try to use image from the specific category first
            if category_images:
                if num < len(category_images):
                    selected_image = category_images[num]
                    print(f"✅ DIRECT: {text} using category image at index {num}")
                else:
                    # Loop within the same category
                    actual_index = num % len(category_images)
                    selected_image = category_images[actual_index]
                    print(f"🔄 CATEGORY LOOP: {text} using category image at index {actual_index} (requested {num})")
            
            # Strategy 2: If no images in specific category, use universal pool
            elif all_images:
                universal_index = num % len(all_images)
                selected_image = all_images[universal_index]
                print(f"🌐 UNIVERSAL LOOP: {text} (no {category} images) using universal image at index {universal_index}")
            
            # Strategy 3: No images available anywhere
            else:
                print(f"⚠️ NO IMAGES AVAILABLE anywhere - keeping placeholder {text}")
                continue  # Keep the placeholder

            # Process the selected image
            if selected_image:
                print(f"🔄 Processing: {text} → {selected_image}")

                # Store placeholder position and size before removing it
                placeholder_left = shape.left
                placeholder_top = shape.top
                placeholder_width = shape.width
                placeholder_height = shape.height
                
                print(f"📍 Placeholder: pos=({placeholder_left}, {placeholder_top}), size={placeholder_width}x{placeholder_height}EMUs")

                # Download if URL, else use local path
                img_file = None
                if isinstance(selected_image, str) and selected_image.startswith("http"):
                    try:
                        resp = requests.get(selected_image, timeout=20)
                        resp.raise_for_status()
                        img_file = BytesIO(resp.content)
                        print(f"✅ Downloaded {len(resp.content)} bytes")
                    except Exception as e:
                        print(f"❌ Download failed: {e}")
                        continue
                else:
                    img_file = selected_image

                try:
                    # Get original image dimensions
                    img_dimensions = get_image_dimensions(img_file if img_file else selected_image)
                    
                    if img_dimensions:
                        img_width, img_height = img_dimensions
                        print(f"📏 Original image: {img_width}x{img_height}px")
                        
                        # Calculate image size with FIXED 540px HEIGHT
                        new_width, new_height = calculate_image_size_for_slide_fixed_height(
                            img_width, img_height, placeholder_width, placeholder_height, fixed_height_px=8000
                        )
                        
                        # Position image at placeholder location (not centered on slide)
                        left = placeholder_left
                        top = placeholder_top
                        
                        # Optional: Center within placeholder bounds if image is smaller
                        if new_width < placeholder_width:
                            left = placeholder_left + (placeholder_width - new_width) // 2
                        if new_height < placeholder_height:
                            top = placeholder_top + (placeholder_height - new_height) // 2
                        
                        # Ensure image doesn't go off-slide bounds
                        left = max(0, min(left, slide_width - new_width))
                        top = max(0, min(top, slide_height - new_height))
                        
                        print(f"🎯 Final placement: {new_width}x{new_height}EMUs at ({left}, {top})")
                        
                    else:
                        # Fallback: use placeholder dimensions and position
                        print("⚠️ Could not detect image dimensions - using placeholder size")
                        left, top, new_width, new_height = placeholder_left, placeholder_top, placeholder_width, placeholder_height

                    # Remove placeholder and add image at its exact position
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    slide.shapes.add_picture(img_file, left, top, new_width, new_height)
                    replaced_any = True
                    print(f"✅ Image inserted at placeholder position with 540px height")

                except Exception as e:
                    print(f"❌ Failed to insert image: {e}")
                    import traceback
                    traceback.print_exc()

        # Mark slides for deletion if no images were added
        if found_placeholder and not replaced_any:
            has_pictures = any(shape.shape_type == 13 for shape in slide.shapes)
            if not has_pictures:
                print(f"🗑️ Marking slide {slide_idx+1} for deletion (no images)")
                slides_to_delete.append(slide)

    prs.save(output_path)
    print(f"💾 Saved presentation: {output_path}")
    return output_path

def get_file_download_url(asset_id: int) -> str:
    """Get the actual downloadable S3 URL for a Monday.com file using the API"""
    query = """
    query($asset_ids: [ID!]!) {
      assets(ids: $asset_ids) {
        id
        name
        public_url
        file_extension
        url
      }
    }
    """
    headers = {
        "Authorization": MONDAY_API_KEY, 
        "Content-Type": "application/json"
    }
    
    try:
        response = requests.post(
            MONDAY_API_URL, 
            json={
                "query": query, 
                "variables": {"asset_ids": [str(asset_id)]}
            }, 
            headers=headers, 
            timeout=15
        )
        response.raise_for_status()
        data = response.json()
        
        if "errors" in data:
            print(f"⚠️ GraphQL errors for asset {asset_id}: {data['errors']}")
            return None
            
        assets = data.get("data", {}).get("assets", [])
        if assets:
            asset = assets[0]
            public_url = asset.get("public_url") or asset.get("url")
            if public_url and public_url != "null":
                return public_url
        
        return None
        
    except Exception as e:
        print(f"⚠️ Failed to get URL for asset {asset_id}: {e}")
        return None

# PDF / DOCX / ZIP Image Extraction
def extract_images_from_pdf(pdf_bytes_or_path):
    images = []
    doc = fitz.open(stream=pdf_bytes_or_path, filetype="pdf") if isinstance(pdf_bytes_or_path, bytes) else fitz.open(pdf_bytes_or_path)
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            images.append(f"https://api.monday.com/v2/file/{base_image['xref']}")  # Placeholder URL
    return images

def extract_images_from_docx(docx_bytes_or_path):
    images = []
    doc = Document(BytesIO(docx_bytes_or_path)) if isinstance(docx_bytes_or_path, bytes) else Document(docx_bytes_or_path)
    for i, rel in enumerate(doc.part._rels.values()):
        if "image" in rel.target_ref:
            images.append(f"https://api.monday.com/v2/file/extracted_{i}")  # Placeholder URL
    return images

def extract_images_from_zip(zip_bytes_or_path):
    images = []
    zip_file = ZipFile(BytesIO(zip_bytes_or_path)) if isinstance(zip_bytes_or_path, bytes) else ZipFile(zip_bytes_or_path)
    for file_name in zip_file.namelist():
        if file_name.lower().endswith((".jpg", ".jpeg", ".png")):
            images.append(f"https://api.monday.com/v2/file/zip_{file_name}")
        elif file_name.lower().endswith(".pdf"):
            images.extend(extract_images_from_pdf(BytesIO(zip_file.read(file_name))))
        elif file_name.lower().endswith(".docx"):
            images.extend(extract_images_from_docx(BytesIO(zip_file.read(file_name))))
    return images

def map_webhook_to_form(event: dict):
    col = event.get("columnValues", {})

    # Extract image URLs from Monday "files" type column 
    image_urls = []
    if "files" in col and col["files"].get("value"):
        try:
            file_data = json.loads(col["files"]["value"])  # Parse JSON from Monday
            image_urls = [asset["url"] for asset in file_data]  # Collect URLs
        except Exception:
            image_urls = []

    # Extract style information from dropdown column
    styles = []
    # Check different possible dropdown column names for styles
    style_columns = ["dropdown", "dropdown0", "dropdown1", "dropdown2", "style_dropdown"]
    
    for col_name in style_columns:
        if col_name in col and col[col_name].get("chosenValues"):
            chosen_values = col[col_name].get("chosenValues", [])
            styles = [v.get("name", "") for v in chosen_values if v.get("name")]
            if styles:
                print(f"🎨 Found styles in column '{col_name}': {styles}")
                break

    return {
        "9. What is the property type": col.get("dropdown76", {}).get("chosenValues", [{}])[0].get("name"),
        "City": col.get("text8", {}).get("value"),
        "city": col.get("text8", {}).get("value"),
        "Country": col.get("country6", {}).get("countryName"),
        "country": col.get("country6", {}).get("countryName"),
        "11. Space to be designed": ", ".join([v.get("name", "") for v in col.get("dropdown0", {}).get("chosenValues", [])]),
        "What is the area size?": col.get("short_text8fr4spel", {}).get("value"),
        "Which style's do you like": ", ".join(
        [s.get("name") for s in col.get("dropdown", {}).get("chosenValues", [])]
    ) if col.get("dropdown", {}).get("chosenValues") else "",
        "5. How old are you": col.get("status", {}).get("label", {}).get("text"),
        "12. How many people will leave in the space": col.get("text1", {}).get("value"),
        "10. What best describes your situation": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),
        "13. Kids": col.get("text2", {}).get("value"), 
        "14. Do you have any pets": col.get("text_1", {}).get("value"),
        "16. Please describe the scope of work": col.get("text37", {}).get("value"),
        "22. Is there any other information…": col.get("long_text3", {}).get("text"),
        "Can you explain your picture selection?": col.get("short_textot656d98", {}).get("value"),
        "15. What words describe best the mood and feel": col.get("short_text5fonuzuu", {}).get("value"),
        "XXXX":col.get("short_text8fr4spel", {}).get("value"),
    }

def fetch_user_details(email: str):
    try:
        url = f"https://migrate.omrsolutions.com/get_user_details.php?email={email}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        print("Error fetching user details:", e)
    return {}

def get_item_files(item_id: int):
    """Get all files for an item with their public URLs"""
    query = """
    query($item_ids: [ID!]!) {
      items(ids: $item_ids) {
        id
        name
        assets {
          id
          name
          public_url
          file_extension
          url
        }
      }
    }
    """
    headers = {"Authorization": MONDAY_API_KEY, "Content-Type": "application/json"}
    try:
        response = requests.post(
            MONDAY_API_URL, 
            json={
                "query": query, 
                "variables": {"item_ids": [str(item_id)]}
            }, 
            headers=headers, 
            timeout=15
        )
        response.raise_for_status()
        data = response.json()
        
        if "errors" in data:
            print(f"⚠️ GraphQL errors for item {item_id}: {data['errors']}")
            return []
            
    except Exception as e:
        print(f"⚠️ Request/JSON error: {e}")
        return []
        
    items = data.get("data", {}).get("items", [])
    if not items:
        return []
        
    assets = items[0].get("assets", [])
    result = []
    
    for asset in assets:
        public_url = asset.get("public_url") or asset.get("url")
        if public_url and public_url != "null":
            result.append({
                "id": asset["id"], 
                "name": asset["name"], 
                "url": public_url, 
                "ext": asset["file_extension"]
            })
    
    return result

def replace_text_in_ppt(template_path: str, output_path: str, text_map: dict):
    """Enhanced text replacement with better debugging and placeholder matching"""
    print(f"🔄 Replacing text in PPT: {template_path} -> {output_path}")
    prs = Presentation(template_path)
    
    replacements_made = 0
    all_found_text = []  # To debug what placeholders exist
    
    # First pass: collect all text to see what's in the template
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):
                    if run.text.strip():
                        all_found_text.append(f"Slide {slide_idx+1}: '{run.text.strip()}'")
    
    print("📝 All text found in template:")
    for text in all_found_text[:20]:  # Show first 20 to avoid spam
        print(f"  {text}")
    
    # Second pass: actual replacement
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):
                    original_text = run.text
                    
                    for placeholder, value in text_map.items():
                        if not value:  # Skip empty values
                            continue
                        
                        # Method 1: Exact match
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
                            replacements_made += 1
                            print(f"✅ Exact match: '{placeholder}' -> '{value}' in slide {slide_idx + 1}")
                        
                        # Method 2: Case-insensitive match
                        elif placeholder.lower() in run.text.lower():
                            import re
                            pattern = re.compile(re.escape(placeholder), re.IGNORECASE)
                            new_text = pattern.sub(str(value), run.text)
                            if new_text != run.text:
                                run.text = new_text
                                replacements_made += 1
                                print(f"✅ Case-insensitive: '{placeholder}' -> '{value}' in slide {slide_idx + 1}")
                        
    print(f"📝 Total replacements made: {replacements_made}")
    
    # Debug: Show what wasn't replaced
    print("\n🔍 Form data that might not have been used:")
    for key, value in text_map.items():
        if value and "area size" in key.lower():
            print(f"  '{key}': '{value}'")
    
    prs.save(output_path)
    return output_path

def categorize_and_collect_images(event: dict) -> dict:
    """
    Categorize images from Monday.com files based on column source.
    - First tries to fetch a public_url (S3).
    - If no public_url, downloads via Monday API and extracts images locally.
    """
    categorized_images = {
        "floor_plans": [],
        "elevation_drawings": [],
        "existing_pictures": [],
        "inspiration_images": []
    }
    
    col_vals = event.get("columnValues", {})
    print(f"🔍 Found {len(col_vals)} column values to process")
    
    # Map file columns to categories
    column_category_map = {
        "files": "floor_plans",
        "fileb3p8t108": "elevation_drawings", 
        "fileh7us51cr": "existing_pictures",
        "files3": "inspiration_images"
    }
    
    for column_name, category in column_category_map.items():
        print(f"🔍 Checking column: {column_name} -> {category}")
        
        if column_name not in col_vals:
            print(f"⚠️ Column {column_name} not found in webhook data")
            continue

        file_data = col_vals[column_name]
        if not (isinstance(file_data, dict) and "files" in file_data):
            print(f"⚠️ Unexpected structure in {column_name}: {file_data}")
            continue

        files_list = file_data["files"]
        print(f"📁 Found {len(files_list)} files in {column_name}")

        for file_info in files_list:
            if not isinstance(file_info, dict):
                print(f"⚠️ File info is not a dict: {file_info}")
                continue

            asset_id = file_info.get("assetId") or file_info.get("id")
            filename = file_info.get("name", "")
            file_ext = file_info.get("extension", "").lower()

            if not asset_id or not filename:
                print(f"⚠️ Missing asset_id or filename: {file_info}")
                continue

            print(f"🔍 Processing file: {filename} (ID: {asset_id}, Ext: {file_ext})")

            try:
                # 1. Try to get a direct public URL
                public_url = get_file_download_url(asset_id)

                if public_url and ("amazonaws.com" in public_url or "s3" in public_url):
                    categorized_images[category].append(public_url)
                    print(f"✅ Public URL added to {category}: {public_url}")
                    continue

                # 2. If no public URL, download via Monday API
                print(f"⬇️ Downloading file {filename} from Monday API (no public_url)")
                headers = {"Authorization": MONDAY_API_KEY}
                download_url = f"{MONDAY_API_URL}/file/{asset_id}"
                response = requests.get(download_url, headers=headers, timeout=20)
                response.raise_for_status()
                file_bytes = response.content

                # 3. Extract images based on file type
                extracted = []
                if file_ext in ["jpg", "jpeg", "png"]:
                    # Save raw image as a temporary file path
                    temp_path = f"temp_{asset_id}.{file_ext}"
                    with open(temp_path, "wb") as f:
                        f.write(file_bytes)
                    extracted.append(temp_path)

                elif file_ext == "pdf":
                    extracted = extract_images_from_pdf(file_bytes)

                elif file_ext == "docx":
                    extracted = extract_images_from_docx(file_bytes)

                elif file_ext == "zip":
                    extracted = extract_images_from_zip(file_bytes)

                else:
                    print(f"⚠️ Unsupported file type: {filename}")
                
                if extracted:
                    categorized_images[category].extend(extracted)
                    print(f"✅ Extracted {len(extracted)} images from {filename}")

            except Exception as e:
                print(f"⚠️ Failed to process {filename}: {e}")

    # Remove empty categories
    categorized_images = {k: v for k, v in categorized_images.items() if v}

    print("📁 Image categorization complete:")
    for category, urls in categorized_images.items():
        print(f"  {category}: {len(urls)} images")
        for i, url in enumerate(urls, 1):
            print(f"    {i}. {url}")

    return categorized_images

def filter_style_slides_optimized(prs, selected_styles):
    """Filter style slides based on selected styles - placeholder function"""
    # This function should be implemented based on your style filtering logic
    print(f"🎨 Filtering slides for styles: {selected_styles}")
    return prs

@app.post("/monday-webhook")
async def monday_webhook(request: Request):
    body = await request.json()
    print("🚀 Webhook received:", json.dumps(body, indent=2))
    
    if "challenge" in body:
        return JSONResponse(content={"challenge": body["challenge"]})
    
    if "event" in body:
        event = body["event"]
        item_id = event.get("pulseId")
        
        # Skip if already processed
        if is_item_processed(item_id):
            print(f"⏭️ Item {item_id} already processed, skipping...")
            return {
                "status": "skipped",
                "message": f"Item {item_id} already processed",
                "item_id": item_id
            }
        
        print(f"🆕 Processing new item: {item_id}")
        
        # Step 1: Map webhook data
        form_data = map_webhook_to_form(event)
        col_vals = event.get("columnValues", {})
        
        # Step 2: Extract styles
        selected_styles = form_data.get("Which style's do you like", "").split(", ") if form_data.get("Which style's do you like") else []
        selected_styles = [style.strip() for style in selected_styles if style.strip()]
        
        # Step 3: Extract email
        email = None
        if "email" in col_vals:
            email_block = col_vals["email"]
            if isinstance(email_block, dict):
                email = email_block.get("email") or email_block.get("text")
        if not email:
            email = form_data.get("Email") or form_data.get("email") or "krgarav@gmail.com"

        # Step 4: Fetch extra details
        user_details = fetch_user_details(email)
        print(f"🔍 User details response: {user_details}")
  
        if user_details.get("status") == "success":
            qd = user_details.get("data", {}).get("quotationdetails", {})
            print(f"📊 Quotation details: {qd}")
            
            if qd.get("area_size"):
                form_data["Q. Area"] = qd["area_size"]
                print(f"✅ Area size: {qd['area_size']}")  
            if qd.get("project_name"):
                project_name_value = qd["project_name"]
                form_data["Q. Project Name"] = project_name_value
                
            if qd.get("residential_type"):
                form_data["Q.Nature of the project"] = qd["residential_type"]
                print(f"✅ Residential type: {qd['residential_type']}")
        else:
            print(f"⚠️ No data found for user {email} in user_details or API call failed")

        # Step 5: Categorize images
        categorized_images = categorize_and_collect_images(event)

        results = {}

        # Generate PowerPoint presentation
        try:
            # Load template presentation
            prs = Presentation(TEMPLATE_PATH)

            # Filter style slides using optimized direct mapping
            if selected_styles:
                prs = filter_style_slides_optimized(prs, selected_styles)

            # Replace text placeholders
            form_data_for_text = {k: v for k, v in form_data.items() if v}
            print(f"📝 Text replacements (area/size related):")
            for key, value in form_data_for_text.items():
                if "area" in key.lower() or "size" in key.lower():
                    print(f"  '{key}': '{value}'")

            replace_text_in_ppt(TEMPLATE_PATH, OUTPUT_PATH, form_data_for_text)  
            insert_city_image_in_ppt(OUTPUT_PATH, OUTPUT_PATH, form_data.get("city"))
            replace_placeholders_with_images(OUTPUT_PATH, OUTPUT_PATH, categorized_images)
            filter_ppt(OUTPUT_PATH, OUTPUT_PATH, selected_styles)
            
            # Mark item as processed to prevent duplicate processing
            mark_item_as_processed(item_id)
         
            results["full"] = {
                "output_file": OUTPUT_PATH,
                "ppt_type": "full",
                "styles_processed": selected_styles,
                "categorized_images": categorized_images,
                "email": email,
                "project_name": form_data.get("Project Name", "Unknown")
            }
            print(f"✅ Full PPT created: {OUTPUT_PATH}")
           
        except Exception as e:
            print(f"⚠️ Failed to create PowerPoint: {e}")
            import traceback
            traceback.print_exc()
            results["full"] = {"error": str(e)}

        return {"status": "processed", "results": results}

    return {"status": "ok", "message": "Webhook received but no event data"}


