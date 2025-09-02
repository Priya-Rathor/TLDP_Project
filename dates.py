from pptx import Presentation
from pptx.util import Inches
from datetime import datetime, timedelta
import platform

# choose correct format string for day without leading zero
if platform.system() == "Windows":
    DAY_FORMAT = "%#d-%b"
else:
    DAY_FORMAT = "%-d-%b"

def build_mapping(start_date=None):
    if start_date is None:
        start_date = datetime.today()
    
    mapping = {}
    
    # Weekday placeholders (Mon, Tue, ...)
    for i in range(1, 8):  
        date = start_date + timedelta(days=i - 1)
        mapping[f"{{{{day{i}}}}}"] = date.strftime("%a")

    # Date placeholders (1-Sep, 2-Sep, ...)
    for i in range(1, 50):  
        date = start_date + timedelta(days=i - 1)
        mapping[f"{{{{d{i}}}}}"] = date.strftime(DAY_FORMAT)  
    
    return mapping

def replace_text_in_frame(text_frame, mapping):
    replaced = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            new_text = run.text
            for ph, val in mapping.items():
                if ph in new_text:
                    new_text = new_text.replace(ph, val)
            if new_text != run.text:
                run.text = new_text
                replaced += 1
    return replaced

def iter_shapes(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # group shape
            yield from iter_shapes(shp.shapes)

def update_calendar_with_bg(input_file, output_file, image_path, start_date=None):
    prs = Presentation(input_file)
    mapping = build_mapping(start_date)
    total_replaced = 0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide in prs.slides:
        # 🔹 Step 1: Insert full-slide image as background
        pic = slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

        # 🔹 Step 2: Send picture to back so text/calendar stays visible
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # 🔹 Step 3: Replace placeholders
        for shp in iter_shapes(slide.shapes):
            if getattr(shp, "has_table", False):
                for row in shp.table.rows:
                    for cell in row.cells:
                        total_replaced += replace_text_in_frame(cell.text_frame, mapping)
            elif getattr(shp, "has_text_frame", False):
                total_replaced += replace_text_in_frame(shp.text_frame, mapping)

    prs.save(output_file)
    print(f"✅ Done — replaced {total_replaced} placeholders and set new background image. Saved to {output_file}")

if __name__ == "__main__":
    start_date = datetime.today()
    update_calendar_with_bg(
        "dates.pptx",
        "updated_calendar.pptx",
        r"D:\TLDP_Project\image\images (5).jpg",  # background image
        start_date
    )
