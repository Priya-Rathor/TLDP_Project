from fastapi import FastAPI, Request
from email_utils import send_email_with_ppt
from fastapi.responses import JSONResponse
from fastapi.responses import JSONResponse
from pptx import Presentation
import os, re, json, requests
from io import BytesIO
from docx import Document
import fitz  # PyMuPDF
from zipfile import ZipFile
from pptx.util import Inches
from brochure import create_brochure_ppt
from email_utils import send_email_with_ppt
from PIL import Image

app = FastAPI()

TEMPLATE_PATH ="template.pptx"
BTEMPLATE_PATH="btemplate.pptx"
OUTPUT_PATH = "output.pptx"
BOUTPUT_PATH = "Boutput.pptx"
MONDAY_API_KEY = os.getenv("MONDAY_API_KEY", "eyJhbGciOiJIUzI1NiJ9.eyJ0aWQiOjU0NjI5MjM1NywiYWFpIjoxMSwidWlkIjo3NDc3Njk5NywiaWFkIjoiMjAyNS0wOC0wNFQwOTo0MzowNS4wMDBaIiwicGVyIjoibWU6d3JpdGUiLCJhY3RpZCI6MTIxNDMyMDQsInJnbiI6InVzZTEifQ.yYeelRXHOZlaxwYHBAvi6eXRzD2fNn1H-jX-Pd8Ukcw")
MONDAY_API_URL = "https://api.monday.com/v2"
STYLE_FOLDER = "styleGuide"

# ✅ ADD: Track processed items to prevent duplicate emails
PROCESSED_ITEMS = set()
EMAIL_SENT_LOG = "email_sent_log.txt"

def load_processed_items():
    """Load previously processed item IDs from file"""
    global PROCESSED_ITEMS
    if os.path.exists(EMAIL_SENT_LOG):
        try:
            with open(EMAIL_SENT_LOG, 'r') as f:
                PROCESSED_ITEMS = set(line.strip() for line in f if line.strip())
            print(f"📋 Loaded {len(PROCESSED_ITEMS)} previously processed items")
        except Exception as e:
            print(f"⚠️ Error loading processed items: {e}")
            PROCESSED_ITEMS = set()
    else:
        PROCESSED_ITEMS = set()

def mark_item_as_processed(item_id):
    """Mark an item as processed and save to file"""
    global PROCESSED_ITEMS
    PROCESSED_ITEMS.add(str(item_id))
    try:
        with open(EMAIL_SENT_LOG, 'a') as f:
            f.write(f"{item_id}\n")
        print(f"✅ Marked item {item_id} as processed")
    except Exception as e:
        print(f"⚠️ Error saving processed item: {e}")

def is_item_processed(item_id):
    """Check if an item has already been processed"""
    return str(item_id) in PROCESSED_ITEMS

# Load processed items on startup
load_processed_items()

def get_image_dimensions(img_path_or_bytes):
    """
    Get original dimensions of an image.
    Returns (width, height) in pixels or None if failed.
    """
    try:
        if isinstance(img_path_or_bytes, str) and img_path_or_bytes.startswith("http"):
            # Download image to get dimensions
            response = requests.get(img_path_or_bytes, timeout=10)
            response.raise_for_status()
            img_bytes = BytesIO(response.content)
            img = Image.open(img_bytes)
        elif isinstance(img_path_or_bytes, str):
            # Local file path
            img = Image.open(img_path_or_bytes)
        else:
            # BytesIO object
            img = Image.open(img_path_or_bytes)
        
        return img.size  # Returns (width, height)
    except Exception as e:
        print(f"⚠️ Failed to get image dimensions: {e}")
        return None

def calculate_image_size_for_slide(img_width, img_height, placeholder_width, placeholder_height, maintain_aspect=True, max_width_px=600):
    """
    Calculate the best size for an image with a maximum width constraint of 800px.
    
    Args:
        img_width, img_height: Original image dimensions in pixels
        placeholder_width, placeholder_height: Placeholder dimensions in PowerPoint EMUs
        maintain_aspect: Whether to maintain aspect ratio
        max_width_px: Maximum width in pixels (default 800)
    
    Returns:
        (new_width, new_height) in PowerPoint units (EMUs)
    """
    print(f"🔍 Input: img={img_width}x{img_height}px, placeholder={placeholder_width}x{placeholder_height}EMUs, max_width={max_width_px}px")
    
    if not maintain_aspect:
        return placeholder_width, placeholder_height
    
    # Convert pixels to PowerPoint EMUs (English Metric Units)
    # 1 inch = 914400 EMUs, assuming 72 DPI (standard for most images)
    PIXELS_PER_INCH = 72
    EMUS_PER_INCH = 914400
    
    # FORCE the width limit - this is critical
    original_width = img_width
    if img_width > max_width_px:
        # Scale down to max width while maintaining aspect ratio
        scale_factor = max_width_px / img_width
        img_width = max_width_px
        img_height = int(img_height * scale_factor)
        print(f"⚡ FORCED scaling: {original_width}x{img_height//scale_factor}px → {img_width}x{img_height}px (factor: {scale_factor:.3f})")
    else:
        print(f"✅ Image width {img_width}px is within {max_width_px}px limit")
    
    # Convert image dimensions from pixels to EMUs
    img_width_emu = int((img_width / PIXELS_PER_INCH) * EMUS_PER_INCH)
    img_height_emu = int((img_height / PIXELS_PER_INCH) * EMUS_PER_INCH)
    
    print(f"📐 Converted to EMUs: {img_width_emu}x{img_height_emu}")
    
    # Check if we need further scaling to fit placeholder bounds
    width_scale = placeholder_width / img_width_emu if img_width_emu > placeholder_width else 1.0
    height_scale = placeholder_height / img_height_emu if img_height_emu > placeholder_height else 1.0
    
    # Use the smaller scale to ensure image fits within placeholder bounds
    final_scale = min(width_scale, height_scale, 1.0)  # Never scale up beyond original
    
    new_width = int(img_width_emu * final_scale)
    new_height = int(img_height_emu * final_scale)
    
    print(f"📏 Final result: {new_width}x{new_height} EMUs (final_scale: {final_scale:.3f})")
    print(f"📊 Size check: {new_width/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}x{new_height/EMUS_PER_INCH*PIXELS_PER_INCH:.0f}px equivalent")
    
    return new_width, new_height

def get_image_dimensions_enhanced(img_path_or_bytes):
    """
    Enhanced version with better error handling and debugging.
    """
    try:
        print(f"🔍 Getting dimensions for: {type(img_path_or_bytes)} {str(img_path_or_bytes)[:100]}...")
        
        if isinstance(img_path_or_bytes, str) and img_path_or_bytes.startswith("http"):
            # Download image to get dimensions
            print(f"⬇️ Downloading image from URL...")
            response = requests.get(img_path_or_bytes, timeout=15)
            response.raise_for_status()
            img_bytes = BytesIO(response.content)
            img = Image.open(img_bytes)
        elif isinstance(img_path_or_bytes, str):
            # Local file path
            print(f"📁 Opening local file...")
            img = Image.open(img_path_or_bytes)
        else:
            # BytesIO object
            print(f"💾 Opening from BytesIO...")
            img_path_or_bytes.seek(0)  # Reset position
            img = Image.open(img_path_or_bytes)
        
        dimensions = img.size  # Returns (width, height)
        print(f"✅ Image dimensions detected: {dimensions[0]}x{dimensions[1]}px")
        return dimensions
    except Exception as e:
        print(f"❌ Failed to get image dimensions: {e}")
        return None
    
    
    
    
def replace_placeholders_with_images(pptx_path, output_path, categorized_images):
    """
    Replace placeholders like {{Image1}}, {{Layout2}}, {{Elevation1}}, {{Inspiration1}} 
    with categorized images positioned at the placeholder's exact location.
    """
    prs = Presentation(pptx_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"📏 Slide dimensions: {slide_width} x {slide_height}")

    # Build category mapping
    category_map = {
        "Layout": categorized_images.get("floor_plans", []),
        "Elevation": categorized_images.get("elevation_drawings", []),
        "Inspiration": categorized_images.get("inspiration_images", []),
        "Image": categorized_images.get("existing_pictures", []),
    }

    slides_to_delete = []

    for slide_idx, slide in enumerate(prs.slides):
        replaced_any = False
        found_placeholder = False

        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            if shape.has_text_frame:
                text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs).strip()
            else:
                text = shape.text.strip()

            # Match placeholders like {{CategoryN}} but EXCLUDE style placeholders
            match = re.match(r"\{\{(\w+)(\d+)\}\}", text)
            if not match:
                continue

            category, num = match.group(1), int(match.group(2)) - 1
            
            # Skip style placeholders (they're handled separately)
            if category.lower() == "style":
                continue

            found_placeholder = True
            images = category_map.get(category)

            if images and num < len(images):
                img_path = images[num]
                print(f"\n🔄 Processing: {text} → {img_path}")

                # STORE PLACEHOLDER POSITION AND SIZE BEFORE REMOVING IT
                placeholder_left = shape.left
                placeholder_top = shape.top
                placeholder_width = shape.width
                placeholder_height = shape.height
                
                print(f"📍 Placeholder: pos=({placeholder_left}, {placeholder_top}), size={placeholder_width}x{placeholder_height}EMUs")

                # Download if URL, else use local path
                img_file = None
                if isinstance(img_path, str) and img_path.startswith("http"):
                    try:
                        resp = requests.get(img_path, timeout=20)
                        resp.raise_for_status()
                        img_file = BytesIO(resp.content)
                        print(f"✅ Downloaded {len(resp.content)} bytes")
                    except Exception as e:
                        print(f"❌ Download failed: {e}")
                        continue
                else:
                    img_file = img_path

                try:
                    # Get original image dimensions
                    img_dimensions = get_image_dimensions_enhanced(img_file if img_file else img_path)
                    
                    if img_dimensions:
                        img_width, img_height = img_dimensions
                        print(f"📏 Original image: {img_width}x{img_height}px")
                        
                        # Calculate image size to fit within placeholder bounds with 600px max width
                        new_width, new_height = calculate_image_size_for_slide(
                            img_width, img_height, placeholder_width, placeholder_height, max_width_px=600
                        )
                        
                        # POSITION IMAGE AT PLACEHOLDER LOCATION (NOT CENTERED ON SLIDE)
                        left = placeholder_left
                        top = placeholder_top
                        
                        # Optional: Center within placeholder bounds if image is smaller
                        if new_width < placeholder_width:
                            left = placeholder_left + (placeholder_width - new_width) // 2
                        if new_height < placeholder_height:
                            top = placeholder_top + (placeholder_height - new_height) // 2
                        
                        # Ensure image doesn't go off-slide bounds
                        left = max(0, min(left, slide_width - new_width))
                        top = max(0, min(top, slide_height - new_height))
                        
                        print(f"🎯 Final placement: {new_width}x{new_height}EMUs at ({left}, {top})")
                        
                    else:
                        # Fallback: use placeholder dimensions and position
                        print("⚠️ Could not detect image dimensions - using placeholder size")
                        left, top, new_width, new_height = placeholder_left, placeholder_top, placeholder_width, placeholder_height

                    # Remove placeholder and add image at its exact position
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    slide.shapes.add_picture(img_file, left, top, new_width, new_height)
                    replaced_any = True
                    print(f"✅ Image inserted at placeholder position")

                except Exception as e:
                    print(f"❌ Failed to insert image: {e}")
                    import traceback
                    traceback.print_exc()
                    
            else:
                print(f"⚠️ No image available for {text}")
                shape.text = ""  # Clear placeholder if no image

        # Mark slides for deletion if no images were added
        if found_placeholder and not replaced_any:
            has_pictures = any(shape.shape_type == 13 for shape in slide.shapes)
            if not has_pictures:
                print(f"🗑️ Marking slide {slide_idx+1} for deletion (no images)")
                slides_to_delete.append(slide)

    # Delete marked slides
    for slide in slides_to_delete:
        delete_slide(prs, slide)

    prs.save(output_path)
    print(f"💾 Saved presentation: {output_path}")
    return output_path

def get_file_download_url(asset_id: int) -> str:
    """
    Get the actual downloadable S3 URL for a Monday.com file using the API
    """
    query = """
    query($asset_id: [ID!]) {
      assets(ids: $asset_id) {
        id
        name
        public_url
        file_extension
      }
    }
    """
    headers = {
        "Authorization": MONDAY_API_KEY, 
        "Content-Type": "application/json"
    }
    
    try:
        response = requests.post(
            MONDAY_API_URL, 
            json={
                "query": query, 
                "variables": {"asset_id": str(asset_id)}
            }, 
            headers=headers, 
            timeout=15
        )
        response.raise_for_status()
        data = response.json()
        
        assets = data.get("data", {}).get("assets", [])
        if assets and len(assets) > 0:
            public_url = assets[0].get("public_url")
            if public_url:
                print(f"✅ Got S3 URL for asset {asset_id}: {public_url}")
                return public_url
        
        print(f"⚠️ No public URL found for asset {asset_id}")
        return f"https://api.monday.com/v2/file/{asset_id}"  # Fallback
        
    except Exception as e:
        print(f"⚠️ Failed to get URL for asset {asset_id}: {e}")
        return f"https://api.monday.com/v2/file/{asset_id}"  

# ---- PDF / DOCX / ZIP Image Extraction ----

def extract_images_from_pdf(pdf_bytes_or_path):
    images = []
    doc = fitz.open(stream=pdf_bytes_or_path, filetype="pdf") if isinstance(pdf_bytes_or_path, bytes) else fitz.open(pdf_bytes_or_path)
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            images.append(f"https://api.monday.com/v2/file/{base_image['xref']}")  # Placeholder URL
    return images

def extract_images_from_docx(docx_bytes_or_path):
    images = []
    doc = Document(BytesIO(docx_bytes_or_path)) if isinstance(docx_bytes_or_path, bytes) else Document(docx_bytes_or_path)
    for i, rel in enumerate(doc.part._rels.values()):
        if "image" in rel.target_ref:
            images.append(f"https://api.monday.com/v2/file/extracted_{i}")  # Placeholder URL
    return images

def extract_images_from_zip(zip_bytes_or_path):
    images = []
    zip_file = ZipFile(BytesIO(zip_bytes_or_path)) if isinstance(zip_bytes_or_path, bytes) else ZipFile(zip_bytes_or_path)
    for file_name in zip_file.namelist():
        if file_name.lower().endswith((".jpg", ".jpeg", ".png")):
            images.append(f"https://api.monday.com/v2/file/zip_{file_name}")
        elif file_name.lower().endswith(".pdf"):
            images.extend(extract_images_from_pdf(BytesIO(zip_file.read(file_name))))
        elif file_name.lower().endswith(".docx"):
            images.extend(extract_images_from_docx(BytesIO(zip_file.read(file_name))))
    return images

def normalize_style_name(style_name: str) -> str:
    """Normalize style names to safe filenames"""
    clean = style_name.strip().lower()
    clean = re.sub(r"[ \-/]+", "_", clean)      # spaces, slashes, dashes -> underscore
    clean = re.sub(r"[^a-z0-9_]", "", clean)    # remove any other weird chars
    return clean

def get_style_images(selected_styles):
    """
    Match selected style names with images in STYLE_FOLDER.
    Returns a list of image paths for the Style category
    """
    style_images = []
    if not selected_styles:
        return style_images

    for style in selected_styles:
        normalized = normalize_style_name(style)
        for ext in [".jpg", ".jpeg", ".png"]:
            filename = f"{normalized}{ext}"       
            path = os.path.join(STYLE_FOLDER, filename)
            if os.path.exists(path):
                style_images.append(path)
                print(f"🎨 Style matched: {style} -> {path}")
                break
        else:
            print(f"⚠️ Missing image for style: {style} (expected something like {normalized}.png)")

    return style_images

def map_webhook_to_form(event: dict):
    col = event.get("columnValues", {})

    # Extract image URLs from Monday "files" type column 
    image_urls = []
    if "files" in col and col["files"].get("value"):
        try:
            file_data = json.loads(col["files"]["value"])  # Parse JSON from Monday
            image_urls = [asset["url"] for asset in file_data]  # Collect URLs
        except Exception:
            image_urls = []

    # Extract style information from dropdown column
    styles = []
    # Check different possible dropdown column names for styles
    style_columns = ["dropdown", "dropdown0", "dropdown1", "dropdown2", "style_dropdown"]
    
    for col_name in style_columns:
        if col_name in col and col[col_name].get("chosenValues"):
            chosen_values = col[col_name].get("chosenValues", [])
            styles = [v.get("name", "") for v in chosen_values if v.get("name")]
            if styles:
                print(f"🎨 Found styles in column '{col_name}': {styles}")
                break

    return {
        # --- Project Information ---
        "Project Name": event.get("pulseName"),
        "**Project Name**": event.get("pulseName"),
        "project name": event.get("pulseName"),

        "What is the nature of your project?": col.get("status1", {}).get("label", {}).get("text"),
        "what is the nature of your project?": col.get("status1", {}).get("label", {}).get("text"),

        "What is the property type?": col.get("dropdown76", {}).get("chosenValues", [{}])[0].get("name"),
        "what is the property type?": col.get("dropdown76", {}).get("chosenValues", [{}])[0].get("name"),

        "City": col.get("text8", {}).get("value"),
        "city": col.get("text8", {}).get("value"),

        "Country": col.get("country6", {}).get("countryName"),
        "country": col.get("country6", {}).get("countryName"),

        "Space(S) to be designed": ", ".join([v.get("name", "") for v in col.get("dropdown0", {}).get("chosenValues", [])]),
        "space(s) to be designed": ", ".join([v.get("name", "") for v in col.get("dropdown0", {}).get("chosenValues", [])]),

        "What is the area size?": col.get("short_text8fr4spel", {}).get("value"),
        "what is the area size?": col.get("short_text8fr4spel", {}).get("value"),

        # --- Style Information ---
        "Which style(s) do you like?": ", ".join(styles) if styles else "",
        "which style(s) do you like?": ", ".join(styles) if styles else "",
        "selected_styles": styles,  # Raw list for processing

        # --- Client Information ---
        "How old are you?": col.get("status", {}).get("label", {}).get("text"),
        "how old are you?": col.get("status", {}).get("label", {}).get("text"),
        "{{How old are you?}}": col.get("status", {}).get("label", {}).get("text"),
        "{{how old are you?}}": col.get("status", {}).get("label", {}).get("text"),

        "How many people will live in the space?": col.get("text1", {}).get("value"),
        "how many people will live in the space?": col.get("text1", {}).get("value"),
        "{{How many people will live in the space?}}": col.get("text1", {}).get("value"),
        "{{how many people will live in the space?}}": col.get("text1", {}).get("value"),

        "What best describes your situation?": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),
        "what best describes your situation?": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),
        "{{What best describes your situation?}}": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),
        "{{what best describes your situation?}}": col.get("single_selecti4d0sw1", {}).get("label", {}).get("text"),

        "Do you have any pets?": col.get("text_1", {}).get("value"),
        "do you have any pets?": col.get("text_1", {}).get("value"),
        "{{Do you have any pets?}}": col.get("text_1", {}).get("value"),
        "{{do you have any pets?}}": col.get("text_1", {}).get("value"),

        "Please describe the scope of work": col.get("text37", {}).get("value"),
        "please describe the scope of work": col.get("text37", {}).get("value"),
        "{{Please describe the scope of work}}": col.get("text37", {}).get("value"),
        "{{please describe the scope of work}}": col.get("text37", {}).get("value"),

        # --- Other Information ---
        "Is there any other information you'd like us to know?": col.get("long_text3", {}).get("text"),
        "is there any other information you'd like us to know?": col.get("long_text3", {}).get("text"),
        "is there any other information you'd like us to know": col.get("long_text3", {}).get("text"),
        "{{Is there any other information you'd like us to know?}}": col.get("long_text3", {}).get("text"),
        "{{is there any other information you'd like us to know?}}": col.get("long_text3", {}).get("text"),
        "{{is there any other information you'd like us to know}}": col.get("long_text3", {}).get("text"),
        "Is there any other information you'd like us to know": col.get("long_text3", {}).get("text"),
        "other information you'd like us to know": col.get("long_text3", {}).get("text"),
        "other information": col.get("long_text3", {}).get("text"),

        # --- Words describe ---
        "Words Describe the feel for space?": col.get("short_text5fonuzuu", {}).get("value"),
        "words describe the feel for space?": col.get("short_text5fonuzuu", {}).get("value"),
        "words describe the feel for space": col.get("short_text5fonuzuu", {}).get("value"),
        "**Words Describe the feel for space?**": col.get("short_text5fonuzuu", {}).get("value"),
        "**words describe the feel for space?**": col.get("short_text5fonuzuu", {}).get("value"),
        "**{{ Words Describe the feel for space? }}**": col.get("short_text5fonuzuu", {}).get("value"),
        "{{ Words Describe the feel for space? }}": col.get("short_text5fonuzuu", {}).get("value"),
        "{{Words Describe the feel for space?}}": col.get("short_text5fonuzuu", {}).get("value"),
        "{{words describe the feel for space?}}": col.get("short_text5fonuzuu", {}).get("value"),
        "{{words describe the feel for space}}": col.get("short_text5fonuzuu", {}).get("value"),
        "Words Describe the feel for space": col.get("short_text5fonuzuu", {}).get("value"),
        "words describe the feel for space": col.get("short_text5fonuzuu", {}).get("value"),
    }






def fetch_user_details(email: str):
    try:
        url = f"https://migrate.omrsolutions.com/get_user_details.php?email={email}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return response.json()
    except Exception as e:
        print("Error fetching user details:", e)
    return {}

def get_item_files(item_id: int):
    """
    Get all files for an item with their public URLs
    Fixed GraphQL query syntax
    """
    query = """
    query($item_ids: [ID!]!) {
      items(ids: $item_ids) {
        id
        name
        assets {
          id
          name
          public_url
          file_extension
          url
        }
      }
    }
    """
    headers = {"Authorization": MONDAY_API_KEY, "Content-Type": "application/json"}
    try:
        response = requests.post(
            MONDAY_API_URL, 
            json={
                "query": query, 
                "variables": {"item_ids": [str(item_id)]}  # Changed to array format
            }, 
            headers=headers, 
            timeout=15
        )
        response.raise_for_status()
        data = response.json()
        
        if "errors" in data:
            print(f"⚠️ GraphQL errors for item {item_id}: {data['errors']}")
            return []
            
    except Exception as e:
        print(f"⚠️ Request/JSON error: {e}")
        return []
        
    items = data.get("data", {}).get("items", [])
    if not items:
        return []
        
    assets = items[0].get("assets", [])
    result = []
    
    for asset in assets:
        public_url = asset.get("public_url") or asset.get("url")
        if public_url and public_url != "null":
            result.append({
                "id": asset["id"], 
                "name": asset["name"], 
                "url": public_url, 
                "ext": asset["file_extension"]
            })
    
    return result

def replace_text_in_ppt(template_path: str, output_path: str, text_map: dict):
    """
    Enhanced text replacement with better debugging and placeholder matching
    """
    print(f"🔄 Replacing text in PPT: {template_path} -> {output_path}")
    prs = Presentation(template_path)
    
    replacements_made = 0
    all_found_text = []  # To debug what placeholders exist
    
    # First pass: collect all text to see what's in the template
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):
                    if run.text.strip():
                        all_found_text.append(f"Slide {slide_idx+1}: '{run.text.strip()}'")
    
    print("📝 All text found in template:")
    for text in all_found_text[:20]:  # Show first 20 to avoid spam
        print(f"  {text}")
    
    # Second pass: actual replacement
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(p.runs):
                    original_text = run.text
                    
                    for placeholder, value in text_map.items():
                        if not value:  # Skip empty values
                            continue
                        
                        # Method 1: Exact match
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
                            replacements_made += 1
                            print(f"✅ Exact match: '{placeholder}' -> '{value}' in slide {slide_idx + 1}")
                        
                        # Method 2: Case-insensitive match
                        elif placeholder.lower() in run.text.lower():
                            import re
                            pattern = re.compile(re.escape(placeholder), re.IGNORECASE)
                            new_text = pattern.sub(str(value), run.text)
                            if new_text != run.text:
                                run.text = new_text
                                replacements_made += 1
                                print(f"✅ Case-insensitive: '{placeholder}' -> '{value}' in slide {slide_idx + 1}")
                        
                        # Method 3: Partial match (for debugging)
                        elif any(word in run.text.lower() for word in placeholder.lower().split()):
                            print(f"🔍 Partial match found but not replaced: '{placeholder}' in '{run.text.strip()}'")
    
    print(f"📝 Total replacements made: {replacements_made}")
    
    # Debug: Show what wasn't replaced
    print("\n🔍 Form data that might not have been used:")
    for key, value in text_map.items():
        if value and "area size" in key.lower():
            print(f"  '{key}': '{value}'")
    
    prs.save(output_path)
    return output_path

def categorize_and_collect_images(event: dict) -> dict:
    """
    Categorize images from Monday.com files based on column source.
    - First tries to fetch a public_url (S3).
    - If no public_url, downloads via Monday API and extracts images locally.
    """
    categorized_images = {
        "floor_plans": [],
        "elevation_drawings": [],
        "existing_pictures": [],
        "inspiration_images": []
    }
    
    col_vals = event.get("columnValues", {})
    print(f"🔍 Found {len(col_vals)} column values to process")
    
    # Map file columns to categories
    column_category_map = {
        "files": "floor_plans",
        "fileb3p8t108": "elevation_drawings", 
        "fileh7us51cr": "existing_pictures",
        "files3": "inspiration_images"
    }
    
    for column_name, category in column_category_map.items():
        print(f"🔍 Checking column: {column_name} -> {category}")
        
        if column_name not in col_vals:
            print(f"⚠️ Column {column_name} not found in webhook data")
            continue

        file_data = col_vals[column_name]
        if not (isinstance(file_data, dict) and "files" in file_data):
            print(f"⚠️ Unexpected structure in {column_name}: {file_data}")
            continue

        files_list = file_data["files"]
        print(f"📁 Found {len(files_list)} files in {column_name}")

        for file_info in files_list:
            if not isinstance(file_info, dict):
                print(f"⚠️ File info is not a dict: {file_info}")
                continue

            asset_id = file_info.get("assetId") or file_info.get("id")
            filename = file_info.get("name", "")
            file_ext = file_info.get("extension", "").lower()

            if not asset_id or not filename:
                print(f"⚠️ Missing asset_id or filename: {file_info}")
                continue

            print(f"🔍 Processing file: {filename} (ID: {asset_id}, Ext: {file_ext})")

            try:
                # 1. Try to get a direct public URL
                public_url = get_file_download_url(asset_id)

                if public_url and ("amazonaws.com" in public_url or "s3" in public_url):
                    categorized_images[category].append(public_url)
                    print(f"✅ Public URL added to {category}: {public_url}")
                    continue

                # 2. If no public URL, download via Monday API
                print(f"⬇️ Downloading file {filename} from Monday API (no public_url)")
                headers = {"Authorization": MONDAY_API_KEY}
                download_url = f"{MONDAY_API_URL}/file/{asset_id}"
                response = requests.get(download_url, headers=headers, timeout=20)
                response.raise_for_status()
                file_bytes = response.content

                # 3. Extract images based on file type
                extracted = []
                if file_ext in ["jpg", "jpeg", "png"]:
                    # Save raw image as a temporary file path
                    temp_path = f"temp_{asset_id}.{file_ext}"
                    with open(temp_path, "wb") as f:
                        f.write(file_bytes)
                    extracted.append(temp_path)

                elif file_ext == "pdf":
                    extracted = extract_images_from_pdf(file_bytes)

                elif file_ext == "docx":
                    extracted = extract_images_from_docx(file_bytes)

                elif file_ext == "zip":
                    extracted = extract_images_from_zip(file_bytes)

                else:
                    print(f"⚠️ Unsupported file type: {filename}")
                
                if extracted:
                    categorized_images[category].extend(extracted)
                    print(f"✅ Extracted {len(extracted)} images from {filename}")

            except Exception as e:
                print(f"⚠️ Failed to process {filename}: {e}")

    # Remove empty categories
    categorized_images = {k: v for k, v in categorized_images.items() if v}

    print("📁 Image categorization complete:")
    for category, urls in categorized_images.items():
        print(f"  {category}: {len(urls)} images")
        for i, url in enumerate(urls, 1):
            print(f"    {i}. {url}")

    return categorized_images

def get_file_download_url(asset_id: int) -> str:
    """
    Get the actual downloadable S3 URL for a Monday.com file using the API
    """
    query = """
    query($asset_ids: [ID!]!) {
      assets(ids: $asset_ids) {
        id
        name
        public_url
        file_extension
        url
      }
    }
    """
    headers = {
        "Authorization": MONDAY_API_KEY, 
        "Content-Type": "application/json"
    }
    
    try:
        response = requests.post(
            MONDAY_API_URL, 
            json={
                "query": query, 
                "variables": {"asset_ids": [str(asset_id)]}  # ✅ Correct: list format
            }, 
            headers=headers, 
            timeout=15
        )
        response.raise_for_status()
        data = response.json()
        
        if "errors" in data:
            print(f"⚠️ GraphQL errors for asset {asset_id}: {data['errors']}")
            return None
            
        assets = data.get("data", {}).get("assets", [])
        if assets:
            asset = assets[0]
            public_url = asset.get("public_url") or asset.get("url")
            if public_url and public_url != "null":
                return public_url
        
        return None
        
    except Exception as e:
        print(f"⚠️ Failed to get URL for asset {asset_id}: {e}")
        return None

from pptx.oxml.ns import qn

def delete_slide(prs, slide):
    """
    Safely delete a slide from a Presentation object.
    """
    slide_id = slide.slide_id

    # Find the slide's rel_id
    slide_rel_id = None
    for rel_id, rel in prs.part.rels.items():
        if rel.reltype == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" and rel._target == slide.part:
            slide_rel_id = rel_id
            break

    # Remove from sldIdLst (slideIdList)
    sldIdLst = prs.slides._sldIdLst  
    for sldId in list(sldIdLst):
        if sldId.rId == slide_rel_id:
            sldIdLst.remove(sldId)
            break

    # Drop the relationship if it exists
    if slide_rel_id and slide_rel_id in prs.part.rels:
        prs.part.drop_rel(slide_rel_id)

def replace_style_placeholders(pptx_path, output_path, style_images):
    """
    ✅ UPDATED: Replace style placeholders like {{style1}}, {{style2}}, etc. with style images.
    Style images now FILL THE ENTIRE SLIDE.
    """
    prs = Presentation(pptx_path)
    slides_to_delete = []

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    print(f"🎨 Starting style placeholder replacement with {len(style_images)} style images")
    print(f"📏 Slide dimensions: {slide_width} x {slide_height}")

    for slide_idx, slide in enumerate(prs.slides):
        replaced_any = False
        found_style_placeholder = False

        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            
            # Match style placeholders like {{style1}}, {{style2}}, etc. (case-insensitive)
            match = re.match(r"\{\{style(\d+)\}\}", text, re.IGNORECASE)
            if not match:
                continue

            found_style_placeholder = True
            style_num = int(match.group(1)) - 1  # Convert to 0-based index

            if style_num < len(style_images):
                img_path = style_images[style_num]
                print(f"🎨 Replacing {text} with style image: {img_path} (FULL SLIDE)")

                try:
                    # ✅ REMOVE THE PLACEHOLDER SHAPE FIRST
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    # ✅ ADD IMAGE TO FILL ENTIRE SLIDE (position at 0,0 with full dimensions)
                    slide.shapes.add_picture(img_path, 0, 0, slide_width, slide_height)
                    replaced_any = True
                    print(f"✅ Style image {img_path} added as full-slide background")
                    
                except Exception as e:
                    print(f"⚠️ Failed to insert style image {img_path}: {e}")
            else:
                print(f"⚠️ No style image available for {text} (index {style_num}, available: {len(style_images)})")

        # Mark slide for deletion if it had style placeholders but no replacement and no pictures
        if found_style_placeholder and not replaced_any:
            has_pictures = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = picture
            if not has_pictures:
                print(f"🗑️ Marking slide {slide_idx+1} for deletion (no style images found)")
                slides_to_delete.append(slide)

    # Delete marked slides
    for slide in slides_to_delete:
        delete_slide(prs, slide)

    prs.save(output_path)
    print(f"🎨 Style placeholder replacement complete. Saved to {output_path}")
    return output_path

def replace_placeholders_with_images(pptx_path, output_path, categorized_images):
    """
    ✅ UPDATED: Replace placeholders like {{Image1}}, {{Layout2}}, {{Elevation1}}, {{Inspiration1}} 
    with categorized images using ORIGINAL IMAGE SIZES (scaled to fit slide if needed).
    """
    prs = Presentation(pptx_path)

    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"📏 Slide dimensions: {slide_width} x {slide_height}")

    # ✅ Build category mapping (REMOVED Style from here)
    category_map = {
        "Layout": categorized_images.get("floor_plans", []),
        "Elevation": categorized_images.get("elevation_drawings", []),
        "Inspiration": categorized_images.get("inspiration_images", []),
        "Image": categorized_images.get("existing_pictures", []),
    }

    slides_to_delete = []

    for slide_idx, slide in enumerate(prs.slides):
        replaced_any = False
        found_placeholder = False

        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            if shape.has_text_frame:
                text = "".join(r.text for p in shape.text_frame.paragraphs for r in p.runs).strip()
            else:
                text = shape.text.strip()
                

            # Match placeholders like {{CategoryN}} but EXCLUDE style placeholders
            match = re.match(r"\{\{(\w+)(\d+)\}\}", text)
            if not match:
                continue

            category, num = match.group(1), int(match.group(2)) - 1
            
            # Skip style placeholders (they're handled separately)
            if category.lower() == "style":
                continue

            found_placeholder = True
            images = category_map.get(category)

            if images and num < len(images):
                img_path = images[num]
                print(f"✅ Replacing {text} with {img_path} (using original dimensions)")

                # Download if URL, else use local path
                img_file = None
                if isinstance(img_path, str) and img_path.startswith("http"):
                    try:
                        resp = requests.get(img_path, timeout=20)
                        resp.raise_for_status()
                        img_file = BytesIO(resp.content)
                    except Exception as e:
                        print(f"⚠️ Failed to download {img_path}: {e}")
                        continue
                else:
                    img_file = img_path

                try:
                    # ✅ GET ORIGINAL IMAGE DIMENSIONS
                    img_dimensions = get_image_dimensions(img_file if img_file else img_path)
                    
                    if img_dimensions:
                        img_width, img_height = img_dimensions
                        print(f"📏 Original image size: {img_width}x{img_height}px")
                        
                        # ✅ CALCULATE BEST SIZE TO FIT SLIDE WHILE MAINTAINING ASPECT RATIO
                        new_width, new_height = calculate_image_size_for_slide(
                            img_width, img_height, slide_width, slide_height
                        )
                        
                        # ✅ CENTER THE IMAGE ON THE SLIDE
                        left = (slide_width - new_width) // 2
                        top = (slide_height - new_height) // 2
                        
                        print(f"📏 Calculated size: {new_width}x{new_height}, Position: ({left}, {top})")
                        
                    else:
                        # Fallback to placeholder size if we can't get dimensions
                        print("⚠️ Could not get image dimensions, using placeholder size")
                        left, top, new_width, new_height = shape.left, shape.top, shape.width, shape.height

                    # ✅ REMOVE PLACEHOLDER AND ADD IMAGE
                    sp = shape._element
                    sp.getparent().remove(sp)
                    
                    slide.shapes.add_picture(img_file, left, top, new_width, new_height)
                    replaced_any = True
                    print(f"✅ Image inserted with original proportions")

                except Exception as e:
                    print(f"⚠️ Failed to insert image {img_path}: {e}")
                    
            else:
                print(f"⚠️ No image available for {text}, placeholder removed")
                shape.text = ""  # Clear placeholder if no image

        # ✅ If slide had placeholders but no replacement and no pictures → delete it
        if found_placeholder and not replaced_any:
            has_pictures = any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = picture
            if not has_pictures:
                print(f"🗑️ Marking slide {slide_idx+1} for deletion (no images found)")
                slides_to_delete.append(slide)

    # Delete marked slides
    for slide in slides_to_delete:
        delete_slide(prs, slide)

    prs.save(output_path)
    print(f"💾 Saved updated presentation to {output_path}")
    return output_path


@app.post("/monday-webhook")
async def monday_webhook(request: Request):
    body = await request.json()
    print("🚀 Webhook received:", json.dumps(body, indent=2))
    if "challenge" in body:
        return JSONResponse(content={"challenge": body["challenge"]})
    
    if "event" in body:
        event = body["event"]
        item_id = event.get("pulseId")
        
        # ✅ CHECK: Skip if already processed
        if is_item_processed(item_id):
            print(f"⏭️ Item {item_id} already processed, skipping email...")
            return {
                "status": "skipped",
                "message": f"Item {item_id} already processed",
                "item_id": item_id
            }
        
        print(f"🆕 Processing new item: {item_id}")
        
        # --- Step 1: Map webhook data ---
        form_data = map_webhook_to_form(event)
        col_vals = event.get("columnValues", {})
        
        # --- Step 2: Extract styles ---
        selected_styles = form_data.get("selected_styles", [])
        style_images = get_style_images(selected_styles)
        
        # --- Step 3: Extract email ---
        email = None
        if "email" in col_vals:
            email_block = col_vals["email"]
            if isinstance(email_block, dict):
                email = email_block.get("email") or email_block.get("text")
        if not email:
            email = form_data.get("Email") or form_data.get("email") or "krgarav@gmail.com"

        # --- Step 4: Fetch extra details ---
        user_details = fetch_user_details(email)
        print(f"🔍 User details response: {user_details}")  # Debug line

        if user_details.get("status") == "success":
            qd = user_details.get("data", {}).get("quotationdetails", {})
            print(f"📊 Quotation details: {qd}")  # Debug line
            
            if qd.get("area_size"):
                area_size_value = qd["area_size"]
                print(f"✅ Area size found: '{area_size_value}'")  # Debug line
                
                # Add multiple variations to catch different placeholder formats
                form_data["What is the area size?"] = area_size_value
                form_data["what is the area size?"] = area_size_value
                form_data["{{What is the area size?}}"] = area_size_value
                form_data["{{what is the area size?}}"] = area_size_value
                form_data["Area Size"] = area_size_value
                form_data["{{Area Size}}"] = area_size_value
                form_data["area size"] = area_size_value
                form_data["{{area size}}"] = area_size_value
                form_data["AREA SIZE"] = area_size_value
                form_data["{{AREA SIZE}}"] = area_size_value
            else:
                print(f"⚠️ No area_size found in quotation details")
                
            if qd.get("project_name"):
                project_name_value = qd["project_name"]
                print(f"✅ Project name found: '{project_name_value}'")
                form_data["Project Name"] = project_name_value
                form_data["**Project Name**"] = project_name_value
                form_data["project name"] = project_name_value
                form_data["{{Project Name}}"] = project_name_value
                form_data["{{project name}}"] = project_name_value

            if qd.get("residential_type"):
                residential_type_value = qd["residential_type"]
                print(f"✅ Residential type found: '{residential_type_value}'")
                form_data["What is the nature of your project?"] = residential_type_value
                form_data["what is the nature of your project?"] = residential_type_value
                form_data["{{What is the nature of your project?}}"] = residential_type_value
                form_data["{{what is the nature of your project?}}"] = residential_type_value
        else:
            print(f"⚠️ No data found for user {email} in user_details or API call failed")

        # --- Step 5: Categorize images ---
        categorized_images = categorize_and_collect_images(event)

        results = {}

        
        # ======================
        # 3) Generate BROCHURE PPT
        # ======================
        try:
            # Pick best-fit images from categorized_images + styles
            circle_img = (categorized_images.get("existing_pictures") or [None])[0]  
            calendar_bg = (categorized_images.get("existing_pictures") or [None])[0]
            layout_img = (categorized_images.get("floor_plans") or [None])[0]
            layout_bg = (
                (categorized_images.get("existing_pictures") or [None])[1]
                if len(categorized_images.get("existing_pictures", [])) > 1
                else calendar_bg
            )
            extra_images =[]
            extra_images.extend(categorized_images.get("existing_pictures", []))
            extra_images.extend(categorized_images.get("floor_plans", []))
            extra_images.extend(categorized_images.get("elevation_drawings", []))

            brochure_ppt = create_brochure_ppt(
                BTEMPLATE_PATH,
                BOUTPUT_PATH,
                form_data=form_data,
                circle_img=circle_img,
                calendar_bg=calendar_bg,
                layout_img=layout_img,
                layout_bg=layout_bg,
                extra_images=extra_images,
            )

            results["brochure"] = {
                "output_file": brochure_ppt,
                "ppt_type": "brochure",
                "styles_processed": selected_styles,
                "style_images_found": len(style_images),
                "email": email,
                "project_name": form_data.get("Project Name", "Unknown"),
            }

            print(f"✅ Brochure PPT created: {brochure_ppt}")
            # send_email_with_ppt(email, brochure_ppt, form_data)  # optional email sending

        except Exception as e:
            print(f"⚠️ Failed to create brochure PPT: {e}")
            results["brochure"] = {"error": str(e)}

        # ======================
        # 2) Generate FULL PPT
        # ======================
        try:
            form_data_for_text = {k: v for k, v in form_data.items() if k not in ["selected_styles"]}
            
            # Debug: Print what we're about to pass to text replacement
            print(f"📝 Form data for text replacement (area_size related):")
            for key, value in form_data_for_text.items():
                if "area" in key.lower() or "size" in key.lower():
                    print(f"  '{key}': '{value}'")
            
            replace_text_in_ppt(TEMPLATE_PATH, OUTPUT_PATH, form_data_for_text)

            if style_images:
                replace_style_placeholders(OUTPUT_PATH, OUTPUT_PATH, style_images)

            replace_placeholders_with_images(OUTPUT_PATH, OUTPUT_PATH, categorized_images)

            results["full"] = {
                "output_file": OUTPUT_PATH,
                "ppt_type": "full",
                "styles_processed": selected_styles,
                "style_images_found": len(style_images),
                "categorized_images": categorized_images,
                "email": email,
                "project_name": form_data.get("Project Name", "Unknown")
            }
            print(f"✅ Full PPT created: {OUTPUT_PATH}")
           
        except Exception as e:
            print(f"⚠️ Failed to create full PPT: {e}")
            results["full"] = {"error": str(e)}

        # =====================
        # ✅ SEND EMAIL ONLY ONCE
        # =====================
        print("📧 Preparing to send email with PPTs...")
     
        subject = f"Your Project PPTs - {form_data.get('Project Name', 'Untitled Project')}"
        body = (
               "Hello,\n\n"
               "Please find attached the generated brochure and style guide PPTs "
               "for your project.\n\n"
               "Best regards,\nDesign Team"
            )

        file_paths = [OUTPUT_PATH, BOUTPUT_PATH]

        # Filter missing files (send only existing ones)
        file_paths = [f for f in file_paths if os.path.exists(f)]

        # ⚡ Force send only to client
        recipient = "tiwarisonalika6@gmail.com"
        print(f"📦 Files ready for email: {file_paths}")
        
        email_sent = False
        if file_paths:
            try:
                send_email_with_ppt(recipient, subject, body, file_paths)
                print(f"📧 Email sent to {recipient} with {len(file_paths)} attachment(s)")
                email_sent = True
                for  f in file_paths:
                    if os.path.exists(f):
                        os.remove(f)
                        print(f"🗑️ Deleted temporary file: {f}")
            except Exception as e:
                print(f"⚠️ Failed to send email: {e}")
                email_sent = False
        else:
            print("⚠️ No PPT files found to attach in email")

        # ✅ MARK AS PROCESSED (regardless of email success)
        mark_item_as_processed(item_id)
        
        return {
            "status": "success",
            "message": "Both PPTs processed",
            "results": results,
            "item_id": item_id,
            "email_sent": email_sent,
            "processed_timestamp": str(os.path.getmtime(EMAIL_SENT_LOG) if os.path.exists(EMAIL_SENT_LOG) else "N/A")
        }

    return {"status": "ok", "message": "Webhook received but no event data"}


    