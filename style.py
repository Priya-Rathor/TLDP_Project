from pptx import Presentation

STYLE_TO_SLIDE_MAP = {
    "art deco": 10,
    "asian zen": 11,
    "coastal": 12,
    "contemporary": 13,
    "country": 14,
    "eclectic": 15,
    "industrial": 16,
    "mid-century": 17,
    "minimalist": 18,
    "modern": 19,
    "rustic": 20,
    "scandinavian": 21,
    "shabby chic": 22,
    "traditional": 23,
    "transitional": 24,
    "tropical": 25,
    "urban": 26
}

def filter_ppt(input_ppt, output_ppt, selected_styles):
    prs = Presentation(input_ppt)

    # Normalize selected styles to lowercase
    selected_styles = [s.lower() for s in selected_styles]

    # Find slide indexes to remove
    slides_to_remove = []
    for style, slide_number in STYLE_TO_SLIDE_MAP.items():
        if style not in selected_styles:
            slide_index = slide_number - 1  # adjust for 0-based index
            if 0 <= slide_index < len(prs.slides):
                slides_to_remove.append(slide_index)

    # Remove in reverse order (important to avoid index errors)
    for i in sorted(slides_to_remove, reverse=True):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    prs.save(output_ppt)
    print(f"✅ Saved filtered PPT with only selected styles {selected_styles}: {output_ppt}")



