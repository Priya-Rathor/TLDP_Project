from pptx import Presentation
from pptx.util import Inches
import requests
import os
from PIL import Image
from io import BytesIO

def save_image(url, city_name, width=489, height=540):
    headers = {"User-Agent": "CityImageFetcher/1.0"}
    try:
        img_response = requests.get(url, headers=headers, stream=True, verify=False)
        if img_response.status_code == 200:
            # Open image with PIL
            img = Image.open(BytesIO(img_response.content))
            # Resize
            img = img.resize((width, height), Image.Resampling.LANCZOS)

            filename = f"{city_name.lower()}.jpg"
            img.save(filename, "JPEG")
            print(f"✅ Image of {city_name} saved as {filename} with size {width}x{height}")
            return filename
        else:
            print(f"❌ Failed to download {city_name} image. Status: {img_response.status_code}")
            return None
    except Exception as e:
        print(f"❌ Error downloading {city_name}: {e}")
        return None


def get_city_image(city_name):
    url = "https://en.wikipedia.org/w/api.php"
    params = {
        "action": "query",
        "titles": city_name,
        "prop": "pageimages",
        "format": "json",
        "pithumbsize": 600
    }
    headers = {"User-Agent": "CityImageFetcher/1.0"}

    response = requests.get(url, params=params, headers=headers).json()
    pages = response.get("query", {}).get("pages", {})

    for _, page in pages.items():
        if "thumbnail" in page:
            img_url = page["thumbnail"]["source"]
            print(f"✅ Found Wikipedia thumbnail for {city_name}: {img_url}")
            return save_image(img_url, city_name)

    print(f"❌ No image found for {city_name}.")
    return None


def insert_city_image_in_ppt(ppt_path, output_path, city_name):
    # Get the image
    img_file = get_city_image(city_name)
    if not img_file:
        print("⚠ No image to insert in PPT")
        return

    # Load PPT
    prs = Presentation(ppt_path)

    # Access slide 2 (index 1)
    slide = prs.slides[1]

    # Find placeholder {{cityImage}} and replace
    for shape in slide.shapes:
        if shape.has_text_frame and "{{LocationImage}}" in shape.text:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            # Delete the placeholder
            sp = shape
            sp.element.getparent().remove(sp.element)
            # Insert image in same position
            slide.shapes.add_picture(img_file, left, top, width, height)
            
            print(f"✅ Inserted {city_name} image into slide 2 at placeholder position")

    # Save updated PPT
    prs.save(output_path)
    print(f"✅ Saved updated PPT as {output_path}")
    
    try:
        os.remove(img_file)  # Clean up downloaded image
        print(f"🗑️ Removed temporary image file {img_file}")
    except Exception as e:
        print(f"⚠️ Could not remove temporary image file {img_file}: {e}")    




# child
# styls
# explain pictures 